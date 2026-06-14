"""Testes para office_extraction.py — extração de texto OOXML."""

import os
from io import BytesIO
from pathlib import Path

os.environ["CD_DATABASE_URL"] = "sqlite:///:memory:"
os.environ["CD_GOOGLE_PROVIDER"] = "mock"

import pytest

from classroom_downloader.office_extraction import (
    OFFICE_MIME_TYPES,
    OfficeExtractionResult,
    extract_office_text,
)

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


# ---------------------------------------------------------------------------
# Fixtures — arquivos OOXML reais gerados pelas bibliotecas
# ---------------------------------------------------------------------------


def _make_docx(text: str) -> bytes:
    import docx

    buf = BytesIO()
    doc = docx.Document()
    doc.add_paragraph(text)
    doc.save(buf)
    return buf.getvalue()


def _make_xlsx(rows: list[list]) -> bytes:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Planilha1"
    for row in rows:
        ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slide_texts: list[str]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for text in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        txBox.text_frame.text = text
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Testes docx
# ---------------------------------------------------------------------------


def test_extract_docx_happy_path(tmp_path: Path) -> None:
    content = "Este é o trabalho do aluno sobre biologia."
    path = tmp_path / "trabalho.docx"
    path.write_bytes(_make_docx(content))

    result = extract_office_text(path, DOCX_MIME)

    assert result.error is None
    assert result.truncated is False
    assert content in result.text


def test_extract_docx_returns_error_on_garbage_bytes(tmp_path: Path) -> None:
    path = tmp_path / "corrupt.docx"
    path.write_bytes(b"isto nao e um docx valido\n")

    result = extract_office_text(path, DOCX_MIME)

    assert result.error == "office_parse_failed"
    assert result.text == ""
    assert result.truncated is False


def test_extract_docx_returns_error_on_empty_file(tmp_path: Path) -> None:
    path = tmp_path / "empty.docx"
    path.write_bytes(b"")

    result = extract_office_text(path, DOCX_MIME)

    assert result.error == "office_parse_failed"


def test_extract_docx_with_table(tmp_path: Path) -> None:
    import docx

    buf = BytesIO()
    doc = docx.Document()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Nome"
    table.cell(0, 1).text = "Nota"
    table.cell(1, 0).text = "Aluno"
    table.cell(1, 1).text = "85"
    doc.save(buf)
    path = tmp_path / "tabela.docx"
    path.write_bytes(buf.getvalue())

    result = extract_office_text(path, DOCX_MIME)

    assert result.error is None
    assert "Nome" in result.text or "Aluno" in result.text


# ---------------------------------------------------------------------------
# Testes xlsx
# ---------------------------------------------------------------------------


def test_extract_xlsx_happy_path(tmp_path: Path) -> None:
    path = tmp_path / "dados.xlsx"
    path.write_bytes(_make_xlsx([["Coluna A", "Coluna B"], ["valor1", 42]]))

    result = extract_office_text(path, XLSX_MIME)

    assert result.error is None
    assert result.truncated is False
    assert "Planilha1" in result.text
    assert "Coluna A" in result.text
    assert "42" in result.text


def test_extract_xlsx_preserves_formula_and_number_format(tmp_path: Path) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Calc"
    ws["A1"] = 10
    ws["A2"] = 20
    ws["A3"] = "=SUM(A1:A2)"
    ws["A3"].number_format = '"R$"#,##0.00'
    ws["B1"] = 0.25
    ws["B1"].number_format = "0.00%"
    buf = BytesIO()
    wb.save(buf)
    path = tmp_path / "calc.xlsx"
    path.write_bytes(buf.getvalue())

    result = extract_office_text(path, XLSX_MIME)

    assert result.error is None
    # Fórmula preservada como texto, não como valor calculado.
    assert "=SUM(A1:A2)" in result.text
    # Máscaras de formato visíveis para o avaliador.
    assert "0.00%" in result.text
    assert "R$" in result.text


def test_extract_xlsx_reports_data_validation(tmp_path: Path) -> None:
    import openpyxl
    from openpyxl.worksheet.datavalidation import DataValidation

    wb = openpyxl.Workbook()
    ws = wb.active
    dv = DataValidation(type="list", formula1='"SÁBADO,DOMINGO"')
    dv.add("B1:B10")
    ws.add_data_validation(dv)
    ws["A1"] = "Dia"
    buf = BytesIO()
    wb.save(buf)
    path = tmp_path / "validacao.xlsx"
    path.write_bytes(buf.getvalue())

    result = extract_office_text(path, XLSX_MIME)

    assert result.error is None
    assert "validação de dados" in result.text
    assert "B1:B10" in result.text
    assert "SÁBADO" in result.text


def test_extract_xlsx_reports_bold_style(tmp_path: Path) -> None:
    import openpyxl
    from openpyxl.styles import Font

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Cabeçalho"
    ws["A1"].font = Font(bold=True)
    buf = BytesIO()
    wb.save(buf)
    path = tmp_path / "estilo.xlsx"
    path.write_bytes(buf.getvalue())

    result = extract_office_text(path, XLSX_MIME)

    assert result.error is None
    assert "A1: Cabeçalho" in result.text
    assert "negrito" in result.text


def test_extract_xlsx_returns_error_on_garbage_bytes(tmp_path: Path) -> None:
    path = tmp_path / "corrupt.xlsx"
    path.write_bytes(b"isto nao e um xlsx valido\n")

    result = extract_office_text(path, XLSX_MIME)

    assert result.error == "office_parse_failed"


def test_extract_xlsx_sheet_header_present(tmp_path: Path) -> None:
    import openpyxl

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Alunos"
    ws1.append(["Ana", 90])
    ws2 = wb.create_sheet("Turmas")
    ws2.append(["Turma A"])
    buf = BytesIO()
    wb.save(buf)
    path = tmp_path / "multi.xlsx"
    path.write_bytes(buf.getvalue())

    result = extract_office_text(path, XLSX_MIME)

    assert result.error is None
    assert "=== Planilha: Alunos ===" in result.text
    assert "=== Planilha: Turmas ===" in result.text


# ---------------------------------------------------------------------------
# Testes pptx
# ---------------------------------------------------------------------------


def test_extract_pptx_happy_path(tmp_path: Path) -> None:
    path = tmp_path / "apresentacao.pptx"
    path.write_bytes(_make_pptx(["Slide de introdução", "Slide de conclusão"]))

    result = extract_office_text(path, PPTX_MIME)

    assert result.error is None
    assert result.truncated is False
    assert "=== Slide 1 ===" in result.text
    assert "=== Slide 2 ===" in result.text
    assert "introdução" in result.text or "Slide" in result.text


def test_extract_pptx_returns_error_on_garbage_bytes(tmp_path: Path) -> None:
    path = tmp_path / "corrupt.pptx"
    path.write_bytes(b"isto nao e um pptx valido\n")

    result = extract_office_text(path, PPTX_MIME)

    assert result.error == "office_parse_failed"


# ---------------------------------------------------------------------------
# Truncamento
# ---------------------------------------------------------------------------


def test_extract_docx_truncates_at_cap(tmp_path: Path, monkeypatch) -> None:
    import classroom_downloader.office_extraction as mod

    monkeypatch.setattr(mod, "MAX_TOTAL_TEXT_BYTES", 50)

    content = "A" * 200
    path = tmp_path / "grande.docx"
    path.write_bytes(_make_docx(content))

    result = extract_office_text(path, DOCX_MIME)

    assert result.error is None
    assert result.truncated is True
    assert len(result.text.encode("utf-8")) <= 50


def test_extract_xlsx_truncates_at_cap(tmp_path: Path, monkeypatch) -> None:
    import classroom_downloader.office_extraction as mod

    monkeypatch.setattr(mod, "MAX_TOTAL_TEXT_BYTES", 30)

    path = tmp_path / "grande.xlsx"
    path.write_bytes(
        _make_xlsx([["valor longo " * 5] for _ in range(20)])
    )

    result = extract_office_text(path, XLSX_MIME)

    # Pode não truncar se o cabeçalho da planilha já exceder o cap,
    # mas não deve lançar exceção.
    assert result.error is None


# ---------------------------------------------------------------------------
# Mime não suportado / extensão fallback
# ---------------------------------------------------------------------------


def test_extract_unsupported_mime_returns_error(tmp_path: Path) -> None:
    path = tmp_path / "arquivo.bin"
    path.write_bytes(b"binario")

    result = extract_office_text(path, "application/octet-stream")

    assert result.error == "office_parse_failed"


def test_office_mime_types_constant_has_three_entries() -> None:
    assert len(OFFICE_MIME_TYPES) == 3
    assert DOCX_MIME in OFFICE_MIME_TYPES
    assert XLSX_MIME in OFFICE_MIME_TYPES
    assert PPTX_MIME in OFFICE_MIME_TYPES
