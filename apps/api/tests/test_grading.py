import os
import json
from datetime import UTC, datetime, timedelta
from pathlib import Path
from uuid import uuid4

os.environ["CD_DATABASE_URL"] = "sqlite:///:memory:"
os.environ["CD_GOOGLE_PROVIDER"] = "mock"

from fastapi.testclient import TestClient
from litellm import exceptions as litellm_exceptions
from sqlmodel import Session, select

from classroom_downloader.database import engine
from classroom_downloader.content_extraction import extract_submission_content
from classroom_downloader.google_provider import ClassroomActivity, ClassroomCourse
from classroom_downloader.main import app, provider_dependency
from classroom_downloader.grading_engine import (
    GradingEngine,
    GradingEngineRequest,
    GradingEngineResult,
)
from classroom_downloader.llm_errors import LlmCallError
from classroom_downloader.models import (
    GradingAiAttempt,
    GradingCriterion,
    GradingFileCache,
    GradingJob,
    GradingScrubCache,
    GradingStatus,
    GradingSubmission,
)
from classroom_downloader.settings import get_settings

from grading_helpers import _infer_provider, _seed_infer_job, _text_submission_file


def _enable_litellm_engine(tmp_path, settings) -> None:
    """Point settings at a local single-model catalog with litellm selected.
    conftest restores the settings singleton after each test."""
    cache_path = tmp_path / "model-prices.json"
    overlay_path = tmp_path / "overlay.json"
    cache_path.write_text(
        '{"openai/gpt-5":{"litellm_provider":"openai","mode":"chat","input_cost_per_token":0.000001,"output_cost_per_token":0.000004}}',
        encoding="utf-8",
    )
    overlay_path.write_text(
        '{"schema_version":1,"default_model":"openai/gpt-5","models":{"openai/gpt-5":{"enabled":true,"use_cases":["grading_draft"]}}}',
        encoding="utf-8",
    )
    settings.grading_cache_path = str(tmp_path / "grading")
    settings.grading_engine = "litellm"
    settings.litellm_model = "openai/gpt-5"
    settings.llm_model_catalog_mode = "local_only"
    settings.llm_model_catalog_cache_path = str(cache_path)
    settings.llm_model_overlay_path = str(overlay_path)


def test_grading_health_mock_engine_is_ready() -> None:
    with TestClient(app) as client:
        body = client.get("/api/grading/health").json()

    assert body["engine"] == "mock"
    assert body["ready"] is True
    assert body["status"] == "mock"


def test_grading_health_ready_when_provider_key_present(tmp_path) -> None:
    _enable_litellm_engine(tmp_path, get_settings())  # conftest sets OPENAI_API_KEY
    with TestClient(app) as client:
        body = client.get("/api/grading/health").json()

    assert body["engine"] == "litellm"
    assert body["model"] == "openai/gpt-5"
    assert body["ready"] is True
    assert body["status"] == "ok"
    assert body["missing_keys"] == []


def test_grading_health_reports_missing_provider_key(tmp_path, monkeypatch) -> None:
    _enable_litellm_engine(tmp_path, get_settings())
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    with TestClient(app) as client:
        body = client.get("/api/grading/health").json()

    assert body["ready"] is False
    assert body["status"] == "provider_key_missing"
    assert "OPENAI_API_KEY" in body["missing_keys"]


def test_grading_health_reports_model_not_enabled(tmp_path) -> None:
    settings = get_settings()
    _enable_litellm_engine(tmp_path, settings)
    settings.litellm_model = "openai/not-in-catalog"
    with TestClient(app) as client:
        body = client.get("/api/grading/health").json()

    assert body["ready"] is False
    assert body["status"] == "model_not_enabled"


def test_draft_returns_503_when_provider_key_missing(tmp_path, monkeypatch) -> None:
    _enable_litellm_engine(tmp_path, get_settings())
    monkeypatch.delenv("OPENAI_API_KEY", raising=False)
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
            },
        ).json()
        response = client.post(f"/api/grading/jobs/{job['id']}/draft")

    assert response.status_code == 503
    assert response.json()["detail"]["code"] == "llm_not_configured"
    assert "provider API key" in response.json()["detail"]["message"]


def test_grading_queue_rejects_global_scans() -> None:
    with TestClient(app) as client:
        response = client.get("/api/grading/queue")

    assert response.status_code == 400
    assert "global scans are disabled" in response.json()["detail"]


def test_grading_queue_is_scoped_to_one_assignment() -> None:
    with TestClient(app) as client:
        response = client.get(
            "/api/grading/queue?course_id=course-1&activity_id=activity-1"
        )

    assert response.status_code == 200
    items = response.json()
    assert [item["activity_id"] for item in items] == ["activity-1"]
    cell_diagram = items[0]
    assert cell_diagram["submission_count"] == 2
    assert cell_diagram["status"] in {"ready", "reviewing", "completed"}


def test_grading_job_creation_uses_direct_assignment_lookup() -> None:
    class DirectLookupProvider:
        def get_course(self, course_id: str) -> ClassroomCourse:
            assert course_id == "course-direct"
            return ClassroomCourse(course_id, "Direct Course", None, "ACTIVE")

        def get_activity(self, course_id: str, activity_id: str) -> ClassroomActivity:
            assert course_id == "course-direct"
            assert activity_id == "activity-direct"
            return ClassroomActivity(
                activity_id,
                course_id,
                "Direct Activity",
                "ASSIGNMENT",
                "PUBLISHED",
                None,
            )

        def list_courses(self):
            raise AssertionError("global course scan should not run")

        def list_activities(self, _course_id: str):
            raise AssertionError("course activity scan should not run")

    app.dependency_overrides[provider_dependency] = lambda: DirectLookupProvider()
    try:
        with TestClient(app) as client:
            response = client.post(
                "/api/grading/jobs",
                json={
                    "course_id": "course-direct",
                    "activity_id": "activity-direct",
                    "rubric_mode": "infer",
                    "teacher_loop": "approve",
                },
            )
    finally:
        app.dependency_overrides.pop(provider_dependency, None)

    assert response.status_code == 200
    body = response.json()
    assert body["course_id"] == "course-direct"
    assert body["activity_id"] == "activity-direct"


def test_job_persists_activity_description(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
    with Session(engine) as session:
        row = session.get(GradingJob, job["id"])
        assert row is not None
        assert row.activity_description
        assert "argument" in row.activity_description.lower()


def test_job_without_activity_description_is_none(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
    with Session(engine) as session:
        row = session.get(GradingJob, job["id"])
        assert row is not None
        assert row.activity_description is None


class _CapturingEngine:
    name = "capture"
    model = None

    def __init__(self, criteria):
        self.criteria = criteria
        self.last_request = None

    def grade(self, request):  # pragma: no cover - not used in inference tests
        raise NotImplementedError

    def infer_rubric(self, request):
        self.last_request = request
        return self.criteria


def test_infer_uses_description_only_when_substantial(tmp_path) -> None:
    from classroom_downloader.grading import infer_job_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")
    long_desc = (
        "Write a persuasive essay of at least three paragraphs that opens with a clear "
        "thesis, supports the argument with at least two pieces of textual evidence, "
        "explains how each piece backs the claim, and ends with a conclusion that "
        "restates the central argument and explains why it matters to the reader."
    )
    fake = _CapturingEngine(
        [
            {"name": "Thesis", "weight": 60, "description": None},
            {"name": "Evidence", "weight": 40, "description": None},
        ]
    )
    provider = _infer_provider([_text_submission_file(1, "activity-infer", "some work")])
    with Session(engine) as session:
        job = _seed_infer_job(session, description=long_desc)
        created = infer_job_criteria(session, job, provider, fake)
    assert fake.last_request.description_only is True
    assert fake.last_request.samples == []
    assert {row.name for row in created} == {"Thesis", "Evidence"}


def test_infer_uses_sample_when_description_thin(tmp_path) -> None:
    from classroom_downloader.grading import infer_job_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")
    fake = _CapturingEngine(
        [
            {"name": "Logic", "weight": 50, "description": None},
            {"name": "Style", "weight": 50, "description": None},
        ]
    )
    files = [
        _text_submission_file(1, "activity-infer", "def add(a, b):\n    return a + b\n"),
        _text_submission_file(2, "activity-infer", "result = 6 * 7\nprint(result)\n"),
    ]
    provider = _infer_provider(files)
    with Session(engine) as session:
        job = _seed_infer_job(session, description="Short.")
        created = infer_job_criteria(session, job, provider, fake)
    assert fake.last_request.description_only is False
    assert len(fake.last_request.samples) == 2
    assert {row.name for row in created} == {"Logic", "Style"}


def test_infer_caps_sample_at_configured_size(tmp_path) -> None:
    from classroom_downloader.grading import infer_job_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")
    get_settings().rubric_infer_sample_size = 2
    fake = _CapturingEngine([{"name": "A", "weight": 100, "description": None}])
    files = [
        _text_submission_file(i, "activity-infer", f"line {i}\nvalue = {i}\n")
        for i in range(5)
    ]
    provider = _infer_provider(files)
    with Session(engine) as session:
        job = _seed_infer_job(session, description=None)
        infer_job_criteria(session, job, provider, fake)
    assert len(fake.last_request.samples) == 2


def test_infer_falls_back_to_defaults_when_no_signal(tmp_path) -> None:
    from classroom_downloader.grading import DEFAULT_CRITERIA, infer_job_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")
    fake = _CapturingEngine([{"name": "ShouldNotUse", "weight": 100, "description": None}])
    provider = _infer_provider([])  # no submissions at all
    with Session(engine) as session:
        job = _seed_infer_job(session, description=None)
        created = infer_job_criteria(session, job, provider, fake)
    assert fake.last_request is None  # engine never called
    assert [row.name for row in created] == [c.name for c in DEFAULT_CRITERIA]


def test_inferred_weights_sum_to_100(tmp_path) -> None:
    from classroom_downloader.grading import infer_job_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")
    long_desc = (
        "Design and document a small program. Explain your approach, justify the data "
        "structures you chose, describe how you tested it, and reflect on what you would "
        "improve next time around if you had more time to keep iterating on the work."
    )
    fake = _CapturingEngine(
        [
            {"name": "Design", "weight": 3, "description": None},
            {"name": "Testing", "weight": 3, "description": None},
            {"name": "Reflection", "weight": 4, "description": None},
        ]
    )
    provider = _infer_provider([])
    with Session(engine) as session:
        job = _seed_infer_job(session, description=long_desc)
        created = infer_job_criteria(session, job, provider, fake)
    assert sum(row.weight for row in created) == 100


def test_grade_loop_no_longer_swaps_criteria(tmp_path) -> None:
    from classroom_downloader.grading import _draft_submission, ensure_default_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")

    class _InferReturningEngine:
        name = "infer-return"
        model = None

        def grade(self, request):
            return GradingEngineResult(
                score=90.0,
                confidence=0.9,
                feedback="ok",
                flags=[],
                criterion_notes=[],
                inferred_criteria=[{"name": "Sneaky", "weight": 100, "description": None}],
            )

        def infer_rubric(self, request):  # pragma: no cover
            return []

    provider = _infer_provider(
        [_text_submission_file(1, "activity-infer", "print('hi')\n")]
    )
    with Session(engine) as session:
        job = _seed_infer_job(session, description=None)
        files = provider.list_submission_files("course-infer", ["activity-infer"])
        from classroom_downloader.grading import _submission_for_file

        submission = _submission_for_file(session, job, files[0])
        _draft_submission(session, job, submission, files, provider, _InferReturningEngine())
        session.commit()
        names = {
            row.name
            for row in session.exec(
                select(GradingCriterion).where(GradingCriterion.job_id == job.id)
            ).all()
        }
    assert "Sneaky" not in names


def test_mock_infer_rubric_returns_weighted_criteria() -> None:
    from classroom_downloader.grading_engine import (
        DEFAULT_GRADING_ENGINE,
        RubricInferenceRequest,
    )

    criteria = DEFAULT_GRADING_ENGINE.infer_rubric(
        RubricInferenceRequest(
            job_id="job-1",
            activity_title="Essay Draft",
            activity_description="Write a persuasive essay with a clear thesis.",
            rubric_text=None,
            samples=[],
            description_only=True,
        )
    )
    assert len(criteria) >= 3
    assert sum(int(row["weight"]) for row in criteria) == 100
    assert all(row.get("name") for row in criteria)


def test_build_sample_xml_wraps_and_escapes_samples() -> None:
    from classroom_downloader.litellm_engine import build_sample_xml

    bundle = build_sample_xml(
        [
            {"label": "Aluno A", "source_label": "a.py", "mime_type": "text/x-python", "content": "x < 1 & y"},
            {"label": "Aluno B", "source_label": "b.txt", "mime_type": "text/plain", "content": "ok"},
        ]
    )
    assert bundle.count("<submission") == 2
    assert "x &lt; 1 &amp; y" in bundle
    assert 'label="Aluno A"' in bundle




def test_draft_no_longer_infers_criteria_inline(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
        client.post(f"/api/grading/jobs/{job['id']}/privacy-audit")
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    names = [row["name"] for row in body["criteria"]]
    assert names == ["Compreens\u00e3o", "Evid\u00eancias", "Racioc\u00ednio", "Clareza"]


def test_brief_mode_sends_rubric_text_and_keeps_default_criteria(monkeypatch, tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")

    captured: list[GradingEngineRequest] = []

    class CapturingEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            captured.append(request)
            return GradingEngineResult(score=80, confidence=0.9, feedback="ok", flags=[])

        def infer_rubric(self, request):  # pragma: no cover - brief never infers
            raise AssertionError("brief mode must not infer a rubric")

    with TestClient(app) as client:
        # activity-3 has a substantial description that WOULD drive inference in
        # infer mode; brief must ignore it and keep the default criteria.
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
                "rubric_text": "Priorize a clareza e exemplos.",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: CapturingEngine())
        client.post(f"/api/grading/jobs/{job['id']}/privacy-audit")
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert body["rubric_mode"] == "brief"
    assert [row["name"] for row in body["criteria"]] == [
        "Compreens\u00e3o",
        "Evid\u00eancias",
        "Racioc\u00ednio",
        "Clareza",
    ]
    assert captured
    assert all(request.rubric_mode == "brief" for request in captured)
    assert all(request.rubric_text == "Priorize a clareza e exemplos." for request in captured)
    assert [c["name"] for c in captured[0].criteria] == [
        "Compreens\u00e3o",
        "Evid\u00eancias",
        "Racioc\u00ednio",
        "Clareza",
    ]


def test_remaining_scope_drafts_only_ungraded_classroom_submissions(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
                "scope": "remaining",
            },
        ).json()
        client.post(f"/api/grading/jobs/{job['id']}/privacy-audit")
        drafted = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert drafted["total_submissions"] == 1
    assert [submission["student_name"] for submission in drafted["submissions"]] == ["Bruno Costa"]



def test_create_job_rejects_unknown_rubric_mode() -> None:
    with TestClient(app) as client:
        response = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "totally-made-up",
                "teacher_loop": "approve",
            },
        )

    assert response.status_code == 422
    assert "rubric mode" in response.json()["detail"].lower()


def test_create_job_exposes_visual_submission_consent() -> None:
    with TestClient(app) as client:
        body = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
                "include_visual_submissions": True,
            },
        ).json()
        fetched = client.get(f"/api/grading/jobs/{body['id']}").json()

    assert body["include_visual_submissions"] is True
    assert fetched["include_visual_submissions"] is True


def test_update_criteria_rejects_empty(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
        response = client.patch(
            f"/api/grading/jobs/{job['id']}/criteria",
            json={"criteria": [{"name": "   ", "weight": 0}]},
        )

    assert response.status_code == 422


def test_multi_file_submission_collapses_into_one_card(monkeypatch, tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")

    captured: list[GradingEngineRequest] = []

    class CapturingEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            captured.append(request)
            return GradingEngineResult(score=85, confidence=0.9, feedback="ok", flags=[])

        def infer_rubric(self, request):  # pragma: no cover - brief never infers
            return []

    with TestClient(app) as client:
        # activity-4 has one student (Júlia) who submitted two attachments.
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-4",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: CapturingEngine())
        audit = client.post(f"/api/grading/jobs/{job['id']}/privacy-audit").json()
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    # Two attachments, but one card graded as a set.
    assert audit["total_files"] == 2
    assert body["total_submissions"] == 1
    submission = body["submissions"][0]
    assert submission["student_name"] == "Júlia Rocha"
    assert sorted(f["source_name"] for f in submission["files"]) == ["parte-1.txt", "parte-2.txt"]
    # Exactly one combined grade call carrying both files' content.
    assert len(captured) == 1
    assert "introducao" in captured[0].content
    assert "conclusao" in captured[0].content


def test_draft_job_creates_submissions_criteria_and_cache(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
        response = client.post(f"/api/grading/jobs/{job['id']}/draft")

    assert response.status_code == 200
    body = response.json()
    assert body["status"] == "reviewing"
    assert body["total_submissions"] == 2
    assert len(body["criteria"]) == 4
    assert len(body["cache_files"]) == 2
    assert all(Path(row["content_hash"]).name for row in body["cache_files"])
    assert all(Path(tmp_path / "grading" / body["id"]).exists() for _ in body["cache_files"])


def test_create_job_persists_teacher_criteria(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    criteria = [
        {"name": "Funcionalidade", "weight": 60, "description": "Resolve o problema proposto."},
        {"name": "Clareza", "weight": 40, "description": "Código e explicação são legíveis."},
    ]

    with TestClient(app) as client:
        body = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "structured",
                "teacher_loop": "approve",
                "criteria": criteria,
            },
        ).json()

    assert [
        {
            "name": row["name"],
            "weight": row["weight"],
            "description": row["description"],
        }
        for row in body["criteria"]
    ] == criteria


def test_infer_mode_replaces_defaults_with_ai_criteria(tmp_path) -> None:
    from classroom_downloader.grading import infer_job_criteria

    get_settings().grading_cache_path = str(tmp_path / "grading")
    fake = _CapturingEngine(
        [
            {"name": "Correção", "weight": 70, "description": "Resultado correto."},
            {"name": "Organização", "weight": 30, "description": "Código organizado."},
        ]
    )
    long_desc = (
        "Implemente um programa que resolva a lógica principal proposta e organize o "
        "código em funções claras. Explique brevemente sua abordagem e garanta que a "
        "saída atenda exatamente ao formato pedido no enunciado da atividade entregue."
    )
    provider = _infer_provider([])
    with Session(engine) as session:
        job = _seed_infer_job(session, description=long_desc)
        created = infer_job_criteria(session, job, provider, fake)
    assert [(row.name, row.weight, row.description) for row in created] == [
        ("Correção", 70, "Resultado correto."),
        ("Organização", 30, "Código organizado."),
    ]


def test_defaults_used_when_no_criteria_and_engine_silent(monkeypatch, tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")

    class SilentCriteriaEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            return GradingEngineResult(
                score=82,
                confidence=0.8,
                feedback="Draft without inferred criteria.",
                flags=[],
                criterion_notes=[],
            )

    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: SilentCriteriaEngine())
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert [row["name"] for row in body["criteria"]] == [
        "Compreens\u00e3o",
        "Evid\u00eancias",
        "Racioc\u00ednio",
        "Clareza",
    ]


def test_teacher_loop_off_prepares_rows_without_ai_attempts(monkeypatch, tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")

    class FailingEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            raise AssertionError("off mode must not call the grading engine")

    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "infer",
                "teacher_loop": "off",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: FailingEngine())
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert body["total_submissions"] == 2
    assert all(row["ai_attempt_status"] is None for row in body["submissions"])
    assert all(row["ai_score"] is None for row in body["submissions"])

    with Session(engine) as session:
        attempts = session.exec(
            select(GradingAiAttempt).where(GradingAiAttempt.job_id == job["id"])
        ).all()

    assert attempts == []


def test_teacher_loop_auto_accepts_clean_high_confidence_drafts(monkeypatch, tmp_path) -> None:
    settings = get_settings()
    original_threshold = settings.grading_auto_accept_confidence
    settings.grading_cache_path = str(tmp_path / "grading")
    settings.grading_auto_accept_confidence = 0.85

    class HighConfidenceEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            return GradingEngineResult(
                score=94,
                confidence=0.91,
                feedback="Ready to accept.",
                flags=[],
            )

    try:
        with TestClient(app) as client:
            job = client.post(
                "/api/grading/jobs",
                json={
                    "course_id": "course-2",
                    "activity_id": "activity-3",
                    "rubric_mode": "infer",
                    "teacher_loop": "auto",
                },
            ).json()
            from classroom_downloader import grading

            monkeypatch.setattr(grading, "get_grading_engine", lambda: HighConfidenceEngine())
            body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
    finally:
        settings.grading_auto_accept_confidence = original_threshold

    assert body["status"] == "completed"
    assert body["reviewed_submissions"] == body["total_submissions"]
    assert all(row["reviewed"] is True for row in body["submissions"])
    assert all(row["final_score"] == 94 for row in body["submissions"])


def test_teacher_loop_auto_holds_low_confidence_drafts(monkeypatch, tmp_path) -> None:
    settings = get_settings()
    settings.grading_cache_path = str(tmp_path / "grading")
    settings.grading_auto_accept_confidence = 0.85

    class LowConfidenceEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            return GradingEngineResult(
                score=90,
                confidence=0.70,
                feedback="Below the auto-accept threshold.",
                flags=[],
            )

    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "auto",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: LowConfidenceEngine())
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert body["status"] == "reviewing"
    assert body["reviewed_submissions"] == 0
    assert all(row["reviewed"] is False for row in body["submissions"])
    assert all(row["ai_score"] == 90 for row in body["submissions"])


def test_teacher_loop_auto_ignores_pass1_review_noise(monkeypatch, tmp_path) -> None:
    settings = get_settings()
    settings.grading_cache_path = str(tmp_path / "grading")
    settings.grading_auto_accept_confidence = 0.85

    class FlaggedEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            return GradingEngineResult(
                score=96,
                confidence=0.99,
                feedback="High confidence but flagged for a human.",
                flags=["needs_human"],
            )

    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "auto",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: FlaggedEngine())
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert body["status"] == "completed"
    assert all(row["reviewed"] is True for row in body["submissions"])
    assert all(row["flag"] is None for row in body["submissions"])
    assert all(row["ai_score"] == 96 for row in body["submissions"])


def test_teacher_loop_cowrite_keeps_score_empty_and_feedback_as_reasoning(
    monkeypatch,
    tmp_path,
) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")

    class CowriteEngine(GradingEngine):
        name = "capture"
        model = None

        def grade(self, request: GradingEngineRequest) -> GradingEngineResult:
            return GradingEngineResult(
                score=88,
                confidence=0.8,
                feedback="Reasoning for the teacher to turn into a grade.",
                flags=[],
            )

    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "infer",
                "teacher_loop": "cowrite",
            },
        ).json()
        from classroom_downloader import grading

        monkeypatch.setattr(grading, "get_grading_engine", lambda: CowriteEngine())
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    assert body["status"] == "reviewing"
    assert all(row["ai_score"] is None for row in body["submissions"])
    assert all(row["final_score"] is None for row in body["submissions"])
    assert all("Reasoning" in row["feedback"] for row in body["submissions"])


def test_litellm_engine_attempt_metadata_is_persisted(monkeypatch, tmp_path) -> None:
    settings = get_settings()
    original_settings = {
        "grading_cache_path": settings.grading_cache_path,
        "grading_engine": settings.grading_engine,
        "litellm_model": settings.litellm_model,
        "llm_model_catalog_mode": settings.llm_model_catalog_mode,
        "llm_model_catalog_cache_path": settings.llm_model_catalog_cache_path,
        "llm_model_overlay_path": settings.llm_model_overlay_path,
    }
    cache_path = tmp_path / "model-prices.json"
    overlay_path = tmp_path / "overlay.json"
    cache_path.write_text(
        '{"openai/gpt-5":{"litellm_provider":"openai","mode":"chat","input_cost_per_token":0.000001,"output_cost_per_token":0.000004}}',
        encoding="utf-8",
    )
    overlay_path.write_text(
        '{"schema_version":1,"default_model":"openai/gpt-5","models":{"openai/gpt-5":{"enabled":true,"use_cases":["grading_draft"]}}}',
        encoding="utf-8",
    )

    def fake_completion(**kwargs):
        class Choice:
            message = {
                "content": '{"score": 84, "confidence": 0.8, "feedback": "Solid draft.", "criterion_notes": [], "flags": []}'
            }

        class Response:
            choices = [Choice()]
            usage = {
                "prompt_tokens": 100,
                "completion_tokens": 50,
                "total_tokens": 150,
                "cache_read_input_tokens": 25,
                "cache_creation_input_tokens": 10,
            }

        return Response()

    def fake_completion_cost(**kwargs):
        return 0.1234

    monkeypatch.setattr(
        "classroom_downloader.litellm_engine.litellm.completion", fake_completion
    )
    monkeypatch.setattr(
        "classroom_downloader.grading.litellm.completion_cost", fake_completion_cost
    )

    try:
        settings.grading_cache_path = str(tmp_path / "grading")
        settings.grading_engine = "litellm"
        settings.litellm_model = "openai/gpt-5"
        settings.llm_model_catalog_mode = "local_only"
        settings.llm_model_catalog_cache_path = str(cache_path)
        settings.llm_model_overlay_path = str(overlay_path)

        with TestClient(app) as client:
            job = client.post(
                "/api/grading/jobs",
                json={
                    "course_id": "course-2",
                    "activity_id": "activity-3",
                    "rubric_mode": "brief",
                    "teacher_loop": "approve",
                },
            ).json()
            body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
    finally:
        for key, value in original_settings.items():
            setattr(settings, key, value)

    # activity-3 (course-2) now has 3 submissions: docx, xlsx, pptx.
    n = body["submissions_graded"]
    assert n >= 1
    submission = body["submissions"][0]
    assert submission["ai_engine"] == "litellm"
    assert submission["ai_model"] == "openai/gpt-5"
    assert submission["ai_prompt_tokens"] == 100
    assert submission["ai_completion_tokens"] == 50
    assert submission["ai_token_count"] == 150
    assert submission["ai_cached_prompt_tokens"] == 25
    assert submission["ai_cache_write_tokens"] == 10
    assert submission["ai_cost_cents"] == 12.34
    # Includes the class-level outlier review call in addition to per-submission grading.
    assert body["total_prompt_tokens"] == 100 * (n + 1)
    assert body["total_completion_tokens"] == 50 * (n + 1)
    assert body["total_cached_tokens"] == 25 * (n + 1)
    assert round(body["total_cost_cents"], 2) == round(12.34 * (n + 1), 2)
    assert body["ai_engine"] == "litellm"
    assert body["ai_model"] == "openai/gpt-5"
    assert body["ai_mode"] == settings.grading_batch_mode
    assert body["wall_clock_ms"] is not None


def test_litellm_malformed_response_marks_attempt_failed(
    monkeypatch,
    tmp_path,
) -> None:
    settings = get_settings()
    original_settings = {
        "grading_cache_path": settings.grading_cache_path,
        "grading_engine": settings.grading_engine,
        "litellm_model": settings.litellm_model,
        "llm_model_catalog_mode": settings.llm_model_catalog_mode,
        "llm_model_catalog_cache_path": settings.llm_model_catalog_cache_path,
        "llm_model_overlay_path": settings.llm_model_overlay_path,
    }
    cache_path = tmp_path / "model-prices.json"
    overlay_path = tmp_path / "overlay.json"
    cache_path.write_text(
        '{"openai/gpt-5":{"litellm_provider":"openai","mode":"chat","input_cost_per_token":0.000001,"output_cost_per_token":0.000004}}',
        encoding="utf-8",
    )
    overlay_path.write_text(
        '{"schema_version":1,"default_model":"openai/gpt-5","models":{"openai/gpt-5":{"enabled":true,"use_cases":["grading_draft"]}}}',
        encoding="utf-8",
    )

    def fake_completion(**kwargs):
        class Choice:
            message = {"content": "not-json"}

        class Response:
            choices = [Choice()]
            usage = {"prompt_tokens": 12, "completion_tokens": 3, "total_tokens": 15}

        return Response()

    monkeypatch.setattr(
        "classroom_downloader.litellm_engine.litellm.completion",
        fake_completion,
    )

    try:
        settings.grading_cache_path = str(tmp_path / "grading")
        settings.grading_engine = "litellm"
        settings.litellm_model = "openai/gpt-5"
        settings.llm_model_catalog_mode = "local_only"
        settings.llm_model_catalog_cache_path = str(cache_path)
        settings.llm_model_overlay_path = str(overlay_path)

        with TestClient(app) as client:
            job = client.post(
                "/api/grading/jobs",
                json={
                    "course_id": "course-2",
                    "activity_id": "activity-3",
                    "rubric_mode": "brief",
                    "teacher_loop": "approve",
                },
            ).json()
            body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
    finally:
        for key, value in original_settings.items():
            setattr(settings, key, value)

    submission = body["submissions"][0]
    assert submission["ai_attempt_status"] == "failed"
    assert submission["ai_safe_error"] == "malformed_llm_response"
    assert submission["error"] == "malformed_llm_response"
    assert submission["error_retryable"] is True


def test_litellm_503_marks_attempt_retryable_api_unavailable(
    monkeypatch,
    tmp_path,
) -> None:
    settings = get_settings()
    original_settings = {
        "grading_cache_path": settings.grading_cache_path,
        "grading_engine": settings.grading_engine,
        "litellm_model": settings.litellm_model,
        "llm_model_catalog_mode": settings.llm_model_catalog_mode,
        "llm_model_catalog_cache_path": settings.llm_model_catalog_cache_path,
        "llm_model_overlay_path": settings.llm_model_overlay_path,
    }
    cache_path = tmp_path / "model-prices.json"
    overlay_path = tmp_path / "overlay.json"
    cache_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "fetched_at": "2025-01-01T00:00:00Z",
                "models": [
                    {
                        "litellm_model": "openai/gpt-5",
                        "input_cost_per_token": 0.0,
                        "output_cost_per_token": 0.0,
                        "max_input_tokens": 128000,
                        "max_output_tokens": 4096,
                        "supports_response_schema": True,
                    }
                ],
            }
        ),
        encoding="utf-8",
    )
    overlay_path.write_text(
        '{"schema_version":1,"default_model":"openai/gpt-5","models":{"openai/gpt-5":{"enabled":true,"use_cases":["grading_draft"]}}}',
        encoding="utf-8",
    )

    def fake_completion(**kwargs):
        raise litellm_exceptions.ServiceUnavailableError("down", "openai", "gpt-5")

    monkeypatch.setattr(
        "classroom_downloader.litellm_engine.litellm.completion",
        fake_completion,
    )

    try:
        settings.grading_cache_path = str(tmp_path / "grading")
        settings.grading_engine = "litellm"
        settings.litellm_model = "openai/gpt-5"
        settings.llm_model_catalog_mode = "local_only"
        settings.llm_model_catalog_cache_path = str(cache_path)
        settings.llm_model_overlay_path = str(overlay_path)

        with TestClient(app) as client:
            job = client.post(
                "/api/grading/jobs",
                json={
                    "course_id": "course-2",
                    "activity_id": "activity-3",
                    "rubric_mode": "brief",
                    "teacher_loop": "approve",
                },
            ).json()
            body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
    finally:
        for key, value in original_settings.items():
            setattr(settings, key, value)

    submission = body["submissions"][0]
    assert submission["ai_attempt_status"] == "failed"
    assert submission["ai_safe_error"] == "api_unavailable"
    assert submission["error"] == "api_unavailable"
    assert submission["error_retryable"] is True
    with Session(engine) as session:
        attempt = session.exec(
            select(GradingAiAttempt).where(GradingAiAttempt.submission_id == submission["id"])
        ).first()
    assert attempt is not None
    assert attempt.stage == "grading"
    assert attempt.retryable is True


def test_litellm_missing_score_marks_attempt_failed_in_scored_mode(
    monkeypatch,
    tmp_path,
) -> None:
    settings = get_settings()
    original_settings = {
        "grading_cache_path": settings.grading_cache_path,
        "grading_engine": settings.grading_engine,
        "litellm_model": settings.litellm_model,
        "llm_model_catalog_mode": settings.llm_model_catalog_mode,
        "llm_model_catalog_cache_path": settings.llm_model_catalog_cache_path,
        "llm_model_overlay_path": settings.llm_model_overlay_path,
    }
    cache_path = tmp_path / "model-prices.json"
    overlay_path = tmp_path / "overlay.json"
    cache_path.write_text(
        '{"openai/gpt-5":{"litellm_provider":"openai","mode":"chat","input_cost_per_token":0.000001,"output_cost_per_token":0.000004}}',
        encoding="utf-8",
    )
    overlay_path.write_text(
        '{"schema_version":1,"default_model":"openai/gpt-5","models":{"openai/gpt-5":{"enabled":true,"use_cases":["grading_draft"]}}}',
        encoding="utf-8",
    )

    def fake_completion(**kwargs):
        # Well-formed JSON, but the score is missing in a scored (approve) draft.
        class Choice:
            message = {
                "content": '{"confidence": 0.9, "feedback": "No score.", "criterion_notes": [], "flags": []}'
            }

        class Response:
            choices = [Choice()]
            usage = {"prompt_tokens": 12, "completion_tokens": 3, "total_tokens": 15}

        return Response()

    monkeypatch.setattr(
        "classroom_downloader.litellm_engine.litellm.completion",
        fake_completion,
    )

    try:
        settings.grading_cache_path = str(tmp_path / "grading")
        settings.grading_engine = "litellm"
        settings.litellm_model = "openai/gpt-5"
        settings.llm_model_catalog_mode = "local_only"
        settings.llm_model_catalog_cache_path = str(cache_path)
        settings.llm_model_overlay_path = str(overlay_path)

        with TestClient(app) as client:
            job = client.post(
                "/api/grading/jobs",
                json={
                    "course_id": "course-2",
                    "activity_id": "activity-3",
                    "rubric_mode": "brief",
                    "teacher_loop": "approve",
                },
            ).json()
            body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
    finally:
        for key, value in original_settings.items():
            setattr(settings, key, value)

    submission = body["submissions"][0]
    assert submission["ai_attempt_status"] == "failed"
    assert submission["ai_score"] is None
    assert submission["error"] == "malformed_llm_response"
    assert submission["error_retryable"] is True


def test_unsupported_submission_blocks_normal_ai_draft(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "infer",
                "teacher_loop": "approve",
            },
        ).json()
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    visual_submission = next(row for row in body["submissions"] if row["mime_type"].startswith("image/"))
    assert visual_submission["ai_attempt_status"] == "blocked"
    assert visual_submission["privacy_status"] == "failed"
    assert visual_submission["extraction_status"] == "unsupported"
    assert visual_submission["ai_score"] is None
    assert visual_submission["error"] == "unsupported_visual_submission"


def test_visual_submission_with_consent_is_extracted_scrubbed_cached_and_graded(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
                "include_visual_submissions": True,
            },
        ).json()
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    visual_submission = next(row for row in body["submissions"] if row["mime_type"].startswith("image/"))
    assert visual_submission["ai_attempt_status"] == "completed"
    assert visual_submission["extraction_status"] == "supported"
    assert visual_submission["privacy_status"] == "redacted"
    assert "name_visible" in visual_submission["privacy_flags"]
    assert visual_submission["error"] is None
    assert visual_submission["ai_score"] is not None

    with Session(engine) as session:
        attempts = session.exec(
            select(GradingAiAttempt)
            .where(GradingAiAttempt.submission_id == visual_submission["id"])
            .order_by(GradingAiAttempt.created_at)
        ).all()
        image_cache = session.exec(
            select(GradingFileCache)
            .where(GradingFileCache.job_id == job["id"])
            .where(GradingFileCache.mime_type == "image/png")
        ).one()
        scrub_cache = session.exec(
            select(GradingScrubCache)
            .where(GradingScrubCache.job_id == job["id"])
            .where(GradingScrubCache.content_hash == image_cache.content_hash)
        ).one()

    assert [attempt.stage for attempt in attempts] == ["extraction", "grading", "outlier_review"]
    assert scrub_cache.extraction_status == "supported"
    assert "Transcricao visual mock" in scrub_cache.scrubbed_content
    assert json.loads(scrub_cache.redaction_counts_json)["name_visible"] == 1


def test_retryable_vision_error_does_not_write_scrub_cache_and_surfaces_retryable(
    monkeypatch,
    tmp_path,
) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")

    class ApiUnavailableVisionEngine:
        name = "mock"
        model = None

        def grade(self, request):
            return GradingEngineResult(score=80, confidence=0.9, feedback="ok", flags=[])

        def infer_rubric(self, request):
            return []

        def extract_image(self, request):
            raise LlmCallError("api_unavailable", True, "down")

    from classroom_downloader import grading

    monkeypatch.setattr(grading, "get_grading_engine", lambda: ApiUnavailableVisionEngine())

    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
                "include_visual_submissions": True,
            },
        ).json()
        body = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    visual_submission = next(row for row in body["submissions"] if row["mime_type"].startswith("image/"))
    assert visual_submission["ai_attempt_status"] == "blocked"
    assert visual_submission["error"] == "vision_api_unavailable"
    assert visual_submission["error_retryable"] is True

    with Session(engine) as session:
        attempts = session.exec(
            select(GradingAiAttempt)
            .where(GradingAiAttempt.submission_id == visual_submission["id"])
            .order_by(GradingAiAttempt.created_at)
        ).all()
        image_cache = session.exec(
            select(GradingFileCache)
            .where(GradingFileCache.job_id == job["id"])
            .where(GradingFileCache.mime_type == "image/png")
        ).one()
        scrub_rows = session.exec(
            select(GradingScrubCache)
            .where(GradingScrubCache.job_id == job["id"])
            .where(GradingScrubCache.content_hash == image_cache.content_hash)
        ).all()

    assert [(attempt.stage, attempt.safe_error, attempt.retryable) for attempt in attempts] == [
        ("extraction", "vision_api_unavailable", True),
        ("grading", "vision_api_unavailable", True),
    ]
    assert scrub_rows == []


def test_visual_consent_bypasses_stale_unsupported_visual_scrub_cache(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
            },
        ).json()
        first = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
        first_visual = next(row for row in first["submissions"] if row["mime_type"].startswith("image/"))
        assert first_visual["error"] == "unsupported_visual_submission"

        with Session(engine) as session:
            row = session.get(GradingJob, job["id"])
            assert row is not None
            row.include_visual_submissions = True
            session.add(row)
            session.commit()

        second = client.post(f"/api/grading/jobs/{job['id']}/draft").json()

    second_visual = next(row for row in second["submissions"] if row["mime_type"].startswith("image/"))
    assert second_visual["ai_attempt_status"] == "completed"
    assert second_visual["extraction_status"] == "supported"
    assert second_visual["error"] is None


def test_retry_reuses_cache_before_expiry_and_refetches_after(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
            },
        ).json()
        drafted = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
        submission_id = drafted["submissions"][0]["id"]
        before = len(drafted["cache_files"])

        reused = client.post(
            f"/api/grading/jobs/{job['id']}/submissions/{submission_id}/retry"
        ).json()
        assert len(reused["cache_files"]) == before

        with Session(engine) as session:
            cache = session.exec(
                select(GradingFileCache)
                .where(GradingFileCache.job_id == job["id"])
                .where(GradingFileCache.submission_id == submission_id)
            ).first()
            assert cache is not None
            cache.expires_at = datetime.now(UTC) - timedelta(minutes=1)
            session.add(cache)
            session.commit()

        refetched = client.post(
            f"/api/grading/jobs/{job['id']}/submissions/{submission_id}/retry"
        ).json()

    assert len(refetched["cache_files"]) == before + 1


def test_delete_cache_preserves_results_and_marks_metadata(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "structured",
                "teacher_loop": "approve",
            },
        ).json()
        drafted = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
        cache_path = tmp_path / "grading" / job["id"]
        assert cache_path.exists()

        deleted = client.delete(f"/api/grading/jobs/{job['id']}/cache").json()

    assert len(deleted["submissions"]) == len(drafted["submissions"])
    assert deleted["cache_expires_at"] is None
    assert all(row["deleted_at"] for row in deleted["cache_files"])
    assert not cache_path.exists()


def test_grading_csv_includes_teacher_edits(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "saved",
                "teacher_loop": "approve",
            },
        ).json()
        drafted = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
        submission_id = drafted["submissions"][0]["id"]
        client.post(
            f"/api/grading/jobs/{job['id']}/submissions/{submission_id}/review",
            json={
                "final_score": 88,
                "feedback": "Teacher-edited feedback",
                "reviewed": True,
            },
        )
        response = client.get(f"/api/grading/jobs/{job['id']}/export.csv")

    assert response.status_code == 200
    assert "Teacher-edited feedback" in response.text
    assert "88" in response.text


def test_classroom_links_endpoint_backfills_links_and_posted_state(tmp_path) -> None:
    get_settings().grading_cache_path = str(tmp_path / "grading")
    with TestClient(app) as client:
        job = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-2",
                "activity_id": "activity-3",
                "rubric_mode": "brief",
                "teacher_loop": "approve",
            },
        ).json()
        drafted = client.post(f"/api/grading/jobs/{job['id']}/draft").json()
        submission = drafted["submissions"][0]

        linked = client.post(
            f"/api/grading/jobs/{job['id']}/classroom-links"
        ).json()
        linked_submission = linked["submissions"][0]

        assert linked_submission["classroom_submission_id"] == "export-file-4"
        assert (
            linked_submission["alternate_link"]
            == "https://classroom.google.com/c/course-2/sm/export-file-4/details"
        )
        assert linked_submission["posted_to_classroom"] is False
        assert linked_submission["posted_at"] is None

        posted = client.post(
            f"/api/grading/jobs/{job['id']}/submissions/{submission['id']}/posted",
            json={"posted": True},
        ).json()

        posted_submission = posted["submissions"][0]
        assert posted_submission["posted_to_classroom"] is True
        assert posted_submission["posted_at"] is not None

        reread = client.get(f"/api/grading/jobs/{job['id']}").json()
        assert reread["submissions"][0]["posted_to_classroom"] is True


def _seed_preview_cache(
    tmp_path: Path,
    *,
    mime_type: str,
    source_name: str,
    content: bytes,
) -> tuple[str, str]:
    job_id = f"preview-job-{uuid4()}"
    submission_id = f"preview-submission-{uuid4()}"
    source_file_id = f"file-{uuid4()}"
    cache_dir = tmp_path / job_id
    cache_dir.mkdir(parents=True)
    cached_path = cache_dir / source_name
    cached_path.write_bytes(content)
    with Session(engine) as session:
        session.add(
            GradingJob(
                id=job_id,
                course_id="course-preview",
                course_name="Preview Course",
                activity_id="activity-preview",
                activity_title="Preview Activity",
                rubric_mode="infer",
                teacher_loop="approve",
                total_submissions=1,
                user_email="teacher@example.edu",
            )
        )
        session.add(
            GradingSubmission(
                id=submission_id,
                job_id=job_id,
                source_file_id=source_file_id,
                source_name=source_name,
                mime_type=mime_type,
            )
        )
        session.add(
            GradingFileCache(
                id=f"cache-{uuid4()}",
                job_id=job_id,
                submission_id=submission_id,
                source_file_id=source_file_id,
                source_name=source_name,
                mime_type=mime_type,
                cached_path=str(cached_path),
                content_hash="preview-hash",
                byte_size=len(content),
                expires_at=datetime.now(UTC) + timedelta(minutes=5),
            )
        )
        session.commit()
    return job_id, submission_id


def test_preview_code_file_served_inline_as_text_plain(tmp_path) -> None:
    content = b"def soma(a, b):\n    return a + b\n"
    with TestClient(app) as client:
        job_id, submission_id = _seed_preview_cache(
            tmp_path,
            mime_type="text/x-python",
            source_name="desafios_listas.py",
            content=content,
        )
        response = client.get(f"/api/grading/jobs/{job_id}/submissions/{submission_id}/preview")

    assert response.status_code == 200
    assert response.headers["content-disposition"] == "inline"
    assert response.headers["content-type"] == "text/plain; charset=utf-8"
    assert response.headers["x-content-type-options"] == "nosniff"
    assert response.content == content


def test_preview_binary_still_attachment(tmp_path) -> None:
    content = b"\x00\x01\x02\x03"
    with TestClient(app) as client:
        job_id, submission_id = _seed_preview_cache(
            tmp_path,
            mime_type="application/octet-stream",
            source_name="submission.bin",
            content=content,
        )
        response = client.get(f"/api/grading/jobs/{job_id}/submissions/{submission_id}/preview")

    assert response.status_code == 200
    assert response.headers["content-disposition"].startswith("attachment")
    assert response.headers["x-content-type-options"] == "nosniff"
    assert response.content == content


def test_regrade_creates_new_latest_job() -> None:
    with TestClient(app) as client:
        first = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "structured",
                "teacher_loop": "approve",
                "criteria": [{"name": "Original", "weight": 100}],
            },
        ).json()

        with Session(engine) as session:
            old_job = session.get(GradingJob, first["id"])
            assert old_job is not None
            old_job.status = GradingStatus.completed
            old_job.reviewed_submissions = old_job.total_submissions
            session.add(old_job)
            session.commit()

        second = client.post(
            "/api/grading/jobs",
            json={
                "course_id": "course-1",
                "activity_id": "activity-1",
                "rubric_mode": "brief",
                "teacher_loop": "cowrite",
                "rubric_text": "Use a stricter pass/fail rubric.",
            },
        ).json()
        queue = client.get(
            "/api/grading/queue",
            params={"course_id": "course-1", "activity_id": "activity-1"},
        ).json()
        jobs = client.get("/api/grading/jobs").json()

    assert queue[0]["latest_job_id"] == second["id"]
    latest_activity = next(item for item in jobs if item["activity_id"] == "activity-1")
    assert latest_activity["latest_job_id"] == second["id"]
    assert latest_activity["status"] == "ready"

    with Session(engine) as session:
        assert session.get(GradingJob, first["id"]) is not None
        assert session.get(GradingJob, second["id"]) is not None


# ---------------------------------------------------------------------------
# Lane A — content_extraction: Office OOXML → supported / degraded
# ---------------------------------------------------------------------------


def _make_cache_file(
    tmp_path: Path,
    content: bytes,
    mime_type: str,
    source_name: str = "arquivo",
) -> GradingFileCache:
    """Helper: grava bytes no disco e retorna um GradingFileCache apontando para o arquivo."""
    from io import BytesIO
    path = tmp_path / source_name
    path.write_bytes(content)
    return GradingFileCache(
        id=str(uuid4()),
        job_id="job-test",
        submission_id="sub-test",
        source_file_id="drive-test",
        source_name=source_name,
        mime_type=mime_type,
        cached_path=str(path),
        content_hash="testhash",
        byte_size=len(content),
        expires_at=datetime.now(UTC),
    )


def _make_real_docx() -> bytes:
    import docx
    from io import BytesIO
    buf = BytesIO()
    doc = docx.Document()
    doc.add_paragraph("Trabalho de biologia do aluno.")
    doc.save(buf)
    return buf.getvalue()


def _make_real_xlsx() -> bytes:
    import openpyxl
    from io import BytesIO
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Dados"
    ws.append(["Col", "Valor"])
    ws.append(["A", 1])
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_real_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches
    from io import BytesIO
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    txBox.text_frame.text = "Slide do projeto."
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def test_content_extraction_docx_returns_supported(tmp_path: Path) -> None:
    cache = _make_cache_file(tmp_path, _make_real_docx(), DOCX_MIME, "trabalho.docx")
    result = extract_submission_content(cache)
    assert result.status == "supported"
    assert result.error is None
    assert len(result.text) > 0


def test_content_extraction_xlsx_returns_supported(tmp_path: Path) -> None:
    cache = _make_cache_file(tmp_path, _make_real_xlsx(), XLSX_MIME, "planilha.xlsx")
    result = extract_submission_content(cache)
    assert result.status == "supported"
    assert result.error is None
    assert len(result.text) > 0


def test_content_extraction_pptx_returns_supported(tmp_path: Path) -> None:
    cache = _make_cache_file(tmp_path, _make_real_pptx(), PPTX_MIME, "apresentacao.pptx")
    result = extract_submission_content(cache)
    assert result.status == "supported"
    assert result.error is None
    assert len(result.text) > 0


def test_content_extraction_corrupt_docx_returns_unsupported(tmp_path: Path) -> None:
    cache = _make_cache_file(tmp_path, b"nao e um docx valido", DOCX_MIME, "corrupt.docx")
    result = extract_submission_content(cache)
    assert result.status == "unsupported"
    assert result.error == "office_parse_failed"


def test_content_extraction_js_legacy_mime_returns_supported(tmp_path: Path) -> None:
    # Drive hydrates some .js uploads to the legacy application/x-javascript mime,
    # which used to fall through to unsupported_file_type. It must be supported.
    cache = _make_cache_file(
        tmp_path, b"console.log('oi');\n", "application/x-javascript", "Script.JS"
    )
    result = extract_submission_content(cache)
    assert result.status == "supported"
    assert result.error is None
    assert "console.log" in result.text


def test_content_extraction_js_octet_stream_falls_back_to_extension(tmp_path: Path) -> None:
    # When Drive can't classify the file it reports application/octet-stream;
    # the .js extension fallback (on decodable text) keeps it gradeable.
    cache = _make_cache_file(
        tmp_path, b"const x = 1;\n", "application/octet-stream", "solucao.js"
    )
    result = extract_submission_content(cache)
    assert result.status == "supported"
    assert result.error is None


def test_content_extraction_octet_stream_binary_still_unsupported(tmp_path: Path) -> None:
    # The extension fallback must not rescue genuine binaries: undecodable bytes
    # are rejected as unsupported_binary before the extension check is reached.
    cache = _make_cache_file(
        tmp_path, b"\x00\x01\x02\xff\xfe", "application/octet-stream", "mystery.bin"
    )
    result = extract_submission_content(cache)
    assert result.status == "unsupported"
    assert result.error == "unsupported_binary_submission"


# ---------------------------------------------------------------------------
# Lane B — content_extraction: PDF → pending_vision / unsupported
# ---------------------------------------------------------------------------

MINIMAL_PDF = (
    b"%PDF-1.4\n"
    b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
    b"2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n"
    b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]>>endobj\n"
    b"xref\n0 4\n"
    b"0000000000 65535 f \n"
    b"0000000009 00000 n \n"
    b"0000000058 00000 n \n"
    b"0000000115 00000 n \n"
    b"trailer<</Size 4/Root 1 0 R>>\n"
    b"startxref\n207\n%%EOF\n"
)


def test_content_extraction_pdf_with_visual_flag_returns_pending_vision(tmp_path: Path) -> None:
    cache = _make_cache_file(tmp_path, MINIMAL_PDF, "application/pdf", "relatorio.pdf")
    result = extract_submission_content(cache, allow_visual_pending=True)
    assert result.status == "pending_vision"
    assert result.error is None


def test_content_extraction_pdf_without_visual_flag_returns_unsupported(tmp_path: Path) -> None:
    cache = _make_cache_file(tmp_path, MINIMAL_PDF, "application/pdf", "relatorio.pdf")
    result = extract_submission_content(cache, allow_visual_pending=False)
    assert result.status == "unsupported"
    assert result.error == "unsupported_visual_submission"


def test_content_extraction_image_with_visual_flag_returns_pending_vision(tmp_path: Path) -> None:
    from classroom_downloader.google_provider import MOCK_PNG_BYTES
    cache = _make_cache_file(tmp_path, MOCK_PNG_BYTES, "image/png", "foto.png")
    result = extract_submission_content(cache, allow_visual_pending=True)
    assert result.status == "pending_vision"


# ---------------------------------------------------------------------------
# Lane B — litellm_engine: _build_vision_messages emite parte correta
# ---------------------------------------------------------------------------


def test_build_vision_messages_emits_image_url_for_image() -> None:
    from classroom_downloader.grading_engine import VisionExtractionRequest
    from classroom_downloader.litellm_engine import _build_vision_messages

    req = VisionExtractionRequest(
        job_id="job-1",
        submission_id="sub-1",
        activity_title="Teste",
        source_label="submission.png",
        image_data=b"\x89PNG fake",
        image_mime_type="image/png",
    )
    messages = _build_vision_messages(req)
    user_content = messages[1]["content"]
    types = [part.get("type") for part in user_content]
    assert "image_url" in types
    assert "file" not in types
    image_part = next(p for p in user_content if p.get("type") == "image_url")
    assert "data:image/png;base64," in image_part["image_url"]["url"]


def test_build_vision_messages_emits_file_part_for_pdf() -> None:
    from classroom_downloader.grading_engine import VisionExtractionRequest
    from classroom_downloader.litellm_engine import _build_vision_messages

    req = VisionExtractionRequest(
        job_id="job-1",
        submission_id="sub-1",
        activity_title="Teste",
        source_label="submission.pdf",
        image_data=MINIMAL_PDF,
        image_mime_type="application/pdf",
    )
    messages = _build_vision_messages(req)
    user_content = messages[1]["content"]
    types = [part.get("type") for part in user_content]
    assert "file" in types
    assert "image_url" not in types
    file_part = next(p for p in user_content if p.get("type") == "file")
    assert "data:application/pdf;base64," in file_part["file"]["file_data"]


def test_parse_vision_extraction_result_accepts_pdf_document() -> None:
    from classroom_downloader.litellm_engine import parse_vision_extraction_result

    payload = json.dumps(
        {
            "transcription": "Texto extraido do PDF do aluno.",
            "visual_description": "Documento PDF com tabela de dados.",
            "content_kind": "pdf_document",
            "legibility": "full",
            "pii_observed": [],
        }
    )
    result = parse_vision_extraction_result(payload)
    assert result.content_kind == "pdf_document"
    assert result.transcription == "Texto extraido do PDF do aluno."


