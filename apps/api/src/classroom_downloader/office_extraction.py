"""Extração de texto de arquivos Office OOXML (.docx, .xlsx, .pptx).

Apenas formatos modernos OOXML são suportados. Formatos legados (.doc, .xls,
.ppt) não são tratados. Arquivos que falhem ao parsear retornam
``error="office_parse_failed"`` em vez de lançar exceção — o chamador deve
tratar isso como submissão não suportada.

Para .xlsx, além dos valores das células também são extraídos fórmulas,
máscaras de formato (number_format), validação de dados, formatação condicional
e estilos básicos (negrito/borda/preenchimento) — porque em exercícios de Excel
o artefato avaliado costuma ser a própria formatação, não só os valores.

Imagens e gráficos embutidos nos arquivos Office não são extraídos.
"""

from dataclasses import dataclass
from pathlib import Path

from .observability import get_logger, log_event


logger = get_logger(__name__)


OFFICE_MIME_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

# Mapeamento de extensão de arquivo para tipo Office (fallback quando o mime
# não estiver disponível ou for genérico).
OFFICE_EXTENSIONS = {
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

# Limites de extração — mesma escala que zip_extraction.
MAX_TOTAL_TEXT_BYTES = 500_000
MAX_ROWS_PER_SHEET = 5_000
MAX_CELLS_PER_ROW = 200
# Acima deste tamanho a planilha cai no modo leve (apenas valores), evitando o
# load completo em memória. Exercícios de formatação são pequenos (~KB).
LARGE_XLSX_BYTES = 5 * 1024 * 1024


@dataclass(frozen=True)
class OfficeExtractionResult:
    text: str
    truncated: bool
    error: str | None = None


def extract_office_text(path: Path, mime_type: str) -> OfficeExtractionResult:
    """Extrai texto de um arquivo Office OOXML.

    Retorna ``OfficeExtractionResult`` com ``error="office_parse_failed"`` em
    vez de lançar exceção para entradas inválidas ou corrompidas.
    """
    kind = OFFICE_MIME_TYPES.get(mime_type.lower())
    if kind is None:
        # Tenta pelo sufixo do arquivo como fallback.
        ext = path.suffix.lower()
        resolved_mime = OFFICE_EXTENSIONS.get(ext)
        if resolved_mime:
            kind = OFFICE_MIME_TYPES.get(resolved_mime)
    if kind is None:
        log_event(
            logger,
            "office.extract.unsupported_mime",
            path=str(path),
            mime_type=mime_type,
        )
        return OfficeExtractionResult(text="", truncated=False, error="office_parse_failed")

    try:
        if kind == "docx":
            return _extract_docx(path)
        if kind == "xlsx":
            return _extract_xlsx(path)
        if kind == "pptx":
            return _extract_pptx(path)
    except Exception as exc:  # noqa: BLE001
        log_event(
            logger,
            "office.extract.parse_failed",
            path=str(path),
            mime_type=mime_type,
            kind=kind,
            error=str(exc),
        )
        return OfficeExtractionResult(text="", truncated=False, error="office_parse_failed")

    return OfficeExtractionResult(text="", truncated=False, error="office_parse_failed")


# ---------------------------------------------------------------------------
# Extratores por tipo
# ---------------------------------------------------------------------------

def _extract_docx(path: Path) -> OfficeExtractionResult:
    import docx

    doc = docx.Document(str(path))
    chunks: list[str] = []
    total = 0
    truncated = False

    for block in _docx_blocks(doc):
        encoded = block.encode("utf-8")
        if total + len(encoded) > MAX_TOTAL_TEXT_BYTES:
            truncated = True
            break
        chunks.append(block)
        total += len(encoded)

    text = "\n".join(chunks)
    log_event(
        logger,
        "office.extract.docx",
        path=str(path),
        char_count=len(text),
        truncated=truncated,
    )
    return OfficeExtractionResult(text=text, truncated=truncated)


def _docx_blocks(doc) -> list[str]:
    """Retorna linhas de texto do documento: parágrafos e células de tabelas."""
    blocks: list[str] = []
    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag
        if tag == "p":
            import docx.oxml.ns as _ns
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, doc)
            text = para.text.strip()
            if text:
                blocks.append(text)
        elif tag == "tbl":
            from docx.table import Table
            table = Table(element, doc)
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    blocks.append(line)
    return blocks


def _extract_xlsx(path: Path) -> OfficeExtractionResult:
    # Planilhas grandes caem no modo leve (apenas valores) para evitar o load
    # completo em memória; exercícios de formatação são pequenos.
    if path.stat().st_size > LARGE_XLSX_BYTES:
        return _extract_xlsx_values(path)
    return _extract_xlsx_enriched(path)


def _extract_xlsx_enriched(path: Path) -> OfficeExtractionResult:
    """Extrai conteúdo + estrutura da planilha: fórmulas, máscaras de formato,
    validação de dados, formatação condicional e estilos básicos. Carrega o
    workbook completo (não read_only) para acessar estilos e validações."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), data_only=False)
    chunks: list[str] = []
    total = 0
    truncated = False

    def emit(line: str) -> bool:
        nonlocal total, truncated
        cost = len(line.encode("utf-8")) + 1
        if total + cost > MAX_TOTAL_TEXT_BYTES:
            truncated = True
            return False
        chunks.append(line)
        total += cost
        return True

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        if not emit(f"=== Planilha: {sheet_name} ==="):
            break

        merged = [str(rng) for rng in getattr(ws.merged_cells, "ranges", [])]
        if merged and not emit("[células mescladas] " + ", ".join(merged[:50])):
            break

        rows_read = 0
        for row in ws.iter_rows():
            if rows_read >= MAX_ROWS_PER_SHEET:
                truncated = True
                break
            row_had_content = False
            for cell in row[:MAX_CELLS_PER_ROW]:
                if cell.value is None:
                    continue
                if not emit(_xlsx_cell_line(cell)):
                    break
                row_had_content = True
            if truncated:
                break
            if row_had_content:
                rows_read += 1
        if truncated:
            break

        for line in _xlsx_validations(ws) + _xlsx_conditional(ws):
            if not emit(line):
                break
        if truncated:
            break

    wb.close()
    text = "\n".join(chunks)
    log_event(
        logger,
        "office.extract.xlsx",
        path=str(path),
        char_count=len(text),
        truncated=truncated,
        mode="enriched",
    )
    return OfficeExtractionResult(text=text, truncated=truncated)


def _xlsx_cell_line(cell) -> str:
    """`Coord: valor-ou-fórmula  [fmt: máscara, negrito, borda, preenchimento]`.

    Com data_only=False, células de fórmula retornam o texto da fórmula em
    ``cell.value`` (ex.: ``=SUM(A1:A2)``), preservando exatamente o que o
    exercício avalia."""
    value = cell.value
    text = value if isinstance(value, str) else str(value)
    tags = _xlsx_cell_tags(cell)
    suffix = f"  [{', '.join(tags)}]" if tags else ""
    return f"{cell.coordinate}: {text}{suffix}"


def _xlsx_cell_tags(cell) -> list[str]:
    tags: list[str] = []
    fmt = cell.number_format
    if fmt and fmt != "General":
        tags.append(f"fmt: {fmt}")
    font = cell.font
    if font is not None and font.bold:
        tags.append("negrito")
    border = cell.border
    if border is not None and any(
        getattr(border, side) is not None and getattr(border, side).style
        for side in ("left", "right", "top", "bottom")
    ):
        tags.append("borda")
    fill = cell.fill
    pattern = getattr(fill, "patternType", None) if fill is not None else None
    if pattern not in (None, "none"):
        tags.append("preenchimento")
    return tags


def _xlsx_validations(ws) -> list[str]:
    try:
        dvs = list(ws.data_validations.dataValidation)
    except Exception:  # noqa: BLE001 - validações ausentes/ilegíveis: ignora
        return []
    if not dvs:
        return []
    lines = ["[validação de dados]"]
    for dv in dvs:
        ranges = str(dv.sqref) if dv.sqref is not None else "?"
        dv_type = dv.type or "?"
        formula = dv.formula1 or ""
        lines.append(f"- {ranges} → {dv_type}: {formula}".rstrip())
    return lines


def _xlsx_conditional(ws) -> list[str]:
    entries: list[str] = []
    try:
        cf = ws.conditional_formatting
        for cf_range in cf:
            sqref = getattr(cf_range, "sqref", None)
            for rule in getattr(cf_range, "rules", None) or []:
                entries.append(f"- {sqref} → {getattr(rule, 'type', '?')}")
    except Exception:  # noqa: BLE001 - API instável entre versões: ignora
        return []
    if not entries:
        return []
    return ["[formatação condicional]", *entries]


def _extract_xlsx_values(path: Path) -> OfficeExtractionResult:
    """Modo leve (apenas valores) para planilhas grandes — read_only, sem estilos."""
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    chunks: list[str] = []
    total = 0
    truncated = False

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        header = f"=== Planilha: {sheet_name} ==="
        encoded_header = header.encode("utf-8")
        if total + len(encoded_header) > MAX_TOTAL_TEXT_BYTES:
            truncated = True
            break
        chunks.append(header)
        total += len(encoded_header)

        rows_read = 0
        for row in ws.iter_rows(values_only=True):
            if rows_read >= MAX_ROWS_PER_SHEET:
                truncated = True
                break
            cells = [str(cell) for cell in row[:MAX_CELLS_PER_ROW] if cell is not None]
            if not cells:
                continue
            line = " | ".join(cells)
            encoded = line.encode("utf-8")
            if total + len(encoded) > MAX_TOTAL_TEXT_BYTES:
                truncated = True
                break
            chunks.append(line)
            total += len(encoded)
            rows_read += 1

    wb.close()
    text = "\n".join(chunks)
    log_event(
        logger,
        "office.extract.xlsx",
        path=str(path),
        char_count=len(text),
        truncated=truncated,
        mode="values",
    )
    return OfficeExtractionResult(text=text, truncated=truncated)


def _extract_pptx(path: Path) -> OfficeExtractionResult:
    from pptx import Presentation

    prs = Presentation(str(path))
    chunks: list[str] = []
    total = 0
    truncated = False

    for slide_num, slide in enumerate(prs.slides, start=1):
        header = f"=== Slide {slide_num} ==="
        encoded_header = header.encode("utf-8")
        if total + len(encoded_header) > MAX_TOTAL_TEXT_BYTES:
            truncated = True
            break
        chunks.append(header)
        total += len(encoded_header)

        # Texto dos shapes
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                encoded = text.encode("utf-8")
                if total + len(encoded) > MAX_TOTAL_TEXT_BYTES:
                    truncated = True
                    break
                chunks.append(text)
                total += len(encoded)
            if truncated:
                break

        # Notas do apresentador
        if not truncated and slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            for para in notes_tf.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                encoded = text.encode("utf-8")
                if total + len(encoded) > MAX_TOTAL_TEXT_BYTES:
                    truncated = True
                    break
                chunks.append(text)
                total += len(encoded)

        if truncated:
            break

    text = "\n".join(chunks)
    log_event(
        logger,
        "office.extract.pptx",
        path=str(path),
        char_count=len(text),
        truncated=truncated,
    )
    return OfficeExtractionResult(text=text, truncated=truncated)
