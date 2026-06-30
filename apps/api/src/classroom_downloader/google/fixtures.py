"""Mock fixture bytes and real-corpus loader used by MockGoogleProvider."""
from io import BytesIO
from pathlib import Path


MOCK_PNG_BYTES = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR"
    b"\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00"
    b"\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\xcf\xc0\x00"
    b"\x00\x03\x01\x01\x00\xc9\xfe\x92\xef\x00\x00\x00\x00"
    b"IEND\xaeB`\x82"
)

# Minimal valid PDF (1 página, sem conteúdo visível) para uso como fixture mock.
MOCK_PDF_BYTES = (
    b"%PDF-1.4\n"
    b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
    b"2 0 obj<</Type/Pages/Count 1/Kids[3 0 R]>>endobj\n"
    b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]>>endobj\n"
    b"xref\n0 4\n"
    b"0000000000 65535 f \n"
    b"0000000009 00000 n \n"
    b"0000000058 00000 n \n"
    b"0000000115 00000 n \n"
    b"trailer<</Size 4/Root 1 0 R>>\n"
    b"startxref\n207\n%%EOF\n"
)


def _make_mock_docx() -> bytes:
    """Gera um .docx mínimo real (OOXML) para uso em fixtures mock."""
    import docx as _docx

    buf = BytesIO()
    doc = _docx.Document()
    doc.add_paragraph("Trabalho de dissertacao enviado pelo aluno.")
    doc.add_paragraph("O aluno apresenta argumentos convincentes sobre o tema proposto.")
    doc.save(buf)
    return buf.getvalue()


def _make_mock_xlsx() -> bytes:
    """Gera um .xlsx mínimo real (OOXML) para uso em fixtures mock."""
    import openpyxl as _openpyxl

    wb = _openpyxl.Workbook()
    ws = wb.active
    ws.title = "Dados"
    ws.append(["Nome", "Nota", "Comentario"])
    ws.append(["Aluno Exemplo", 85, "Bom desempenho"])
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_mock_pptx() -> bytes:
    """Gera um .pptx mínimo real (OOXML) para uso em fixtures mock."""
    from pptx import Presentation as _Presentation
    from pptx.util import Inches as _Inches

    prs = _Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(_Inches(1), _Inches(1), _Inches(8), _Inches(4))
    txBox.text_frame.text = "Apresentacao do projeto final do aluno."
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# Gerados uma vez no import; ficam em memória para o ciclo de vida do processo.
MOCK_DOCX_BYTES: bytes = _make_mock_docx()
MOCK_XLSX_BYTES: bytes = _make_mock_xlsx()
MOCK_PPTX_BYTES: bytes = _make_mock_pptx()


# ---------------------------------------------------------------------------
# Real-corpus guarded loader (used by MockGoogleProvider for course-real).
# The loader is best-effort: returns b"" if the corpus or a specific file is
# absent, so the module import never crashes even without the corpus on disk.
# ---------------------------------------------------------------------------

def _find_corpus_root_for_provider() -> "Path | None":
    """Walk up from this file to find the repo root containing test_files/."""
    current = Path(__file__).resolve().parent
    for _ in range(10):
        candidate = current / "test_files"
        if candidate.is_dir():
            return candidate
        parent = current.parent
        if parent == current:
            break
        current = parent
    return None


_CORPUS_ROOT_PROVIDER: "Path | None" = _find_corpus_root_for_provider()


def _corpus_bytes(relative: str) -> bytes:
    """Read bytes from the corpus; return b'' silently if unavailable."""
    if _CORPUS_ROOT_PROVIDER is None:
        return b""
    try:
        return (_CORPUS_ROOT_PROVIDER / relative).read_bytes()
    except OSError:
        return b""
